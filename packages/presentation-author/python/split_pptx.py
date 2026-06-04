#!/usr/bin/env python3
"""Split a multi-slide PPTX into ordered single-slide PPTX files.

Foundation for the per-slide deck Live Preview pipeline (Epic #417 R4): each
single-slide PPTX can then be exported to SVG (export_svg.py) and composed
(compose_slides.py) independently, so slide 1's preview reaches the client
before slide N has finished rendering.

Strategy: for each slide, copy the whole package (preserving masters, layouts,
theme, media) and drop every *other* slide reference from the slide id list.
Copying the package is what keeps each single-slide deck visually identical to
the source — only the referenced slide differs.

Emits a single JSON line on stdout:
  {"success": true,  "slides": [{"index": 1, "pptxPath": "/abs/slide-1.pptx"}, ...]}
  {"success": false, "error": "..."}

Exit code is 0 on success, 1 on failure (so the TS wrapper can degrade).
"""

import argparse
import json
import shutil
import sys
import tempfile
from pathlib import Path


def _emit(payload: dict) -> None:
    sys.stdout.write(json.dumps(payload))
    sys.stdout.flush()


def _keep_only(prs, keep_index: int) -> None:
    """Drop every slide except *keep_index* from the presentation's id list."""
    sld_id_lst = prs.slides._sldIdLst
    sld_ids = list(sld_id_lst)
    for i, sld_id in enumerate(sld_ids):
        if i == keep_index:
            continue
        # Drop the relationship first so the saved package has no dangling r:id,
        # then remove the slide reference from the ordered id list.
        rid = sld_id.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        if rid:
            try:
                prs.part.drop_rel(rid)
            except KeyError:
                pass
        sld_id_lst.remove(sld_id)


def split_pptx(pptx_path: Path, output_dir: Path) -> list[dict]:
    """Write one single-slide PPTX per source slide; return ordered entries."""
    from pptx import Presentation  # imported lazily so a missing dep is reported

    output_dir.mkdir(parents=True, exist_ok=True)

    probe = Presentation(str(pptx_path))
    slide_count = len(probe.slides._sldIdLst)
    if slide_count == 0:
        return []

    slides: list[dict] = []
    with tempfile.TemporaryDirectory(prefix="split_pptx_") as work:
        for index in range(slide_count):
            # Copy the package per slide so each deck keeps masters/layouts/theme.
            tmp_copy = Path(work) / f"copy-{index}.pptx"
            shutil.copyfile(pptx_path, tmp_copy)
            prs = Presentation(str(tmp_copy))
            _keep_only(prs, index)
            out_path = output_dir / f"slide-{index + 1}.pptx"
            prs.save(str(out_path))
            slides.append({"index": index + 1, "pptxPath": str(out_path)})
    return slides


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Split a PPTX into ordered single-slide PPTX files."
    )
    parser.add_argument("pptx_path", help="Path to the input multi-slide .pptx")
    parser.add_argument(
        "--output_dir", required=True, help="Directory for the slide-N.pptx files"
    )
    args = parser.parse_args()

    pptx_path = Path(args.pptx_path)
    output_dir = Path(args.output_dir)

    if not pptx_path.exists():
        _emit({"success": False, "error": f"pptx not found: {pptx_path}"})
        return 1

    try:
        slides = split_pptx(pptx_path, output_dir)
    except ImportError as exc:
        _emit({"success": False, "error": f"python-pptx not installed: {exc}"})
        return 1
    except Exception as exc:  # noqa: BLE001 - degrade with a structured error
        _emit({"success": False, "error": f"{type(exc).__name__}: {exc}"})
        return 1

    if not slides:
        _emit({"success": False, "error": "no slides in presentation"})
        return 1

    _emit({"success": True, "slides": slides})
    return 0


if __name__ == "__main__":
    sys.exit(main())
